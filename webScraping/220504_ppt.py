from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
# 첫번째 슬라이드
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)

title=slide.placeholders[0]
title.text = "Hello, Python!"

subtitle = slide.placeholders[1]
subtitle.text="Python is powerful and fast."

# 두번째 슬라이드
title_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(title_slide_layout)

title_shape = slide.placeholders[0]
title_shape.text = "Python"

body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text="개요"
# 단락 추가
p = tf.add_paragraph()
p.text = "Life is short, You need Python.\n(인생은 짧다. 당신에겐 파이썬이 필요하다.)"
p.level = 1

# 단락 추가(들여쓰기)
p = tf.add_paragraph()
p.text = "1991년에 발표된 인터프리터 방식의 프로그래밍 언어다. 영어와 비슷해서 읽고 쓰기 쉬운 특유의 문법과 그로 인해 전 세계의 개발자들로부터 만들어진 수많은 패키지들 덕분에 사용할 수 있는 용도가 무궁무진해져 2010년대 중반부터 주류 프로그래밍 언어로 당당히 올라서게 되었다."
p.level = 2

prs.save('test.pptx')   